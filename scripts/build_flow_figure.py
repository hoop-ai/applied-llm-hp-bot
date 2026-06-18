"""Build a standalone PORTRAIT flow-diagram figure for the report.

The deck's flow slide is landscape and renders small when fit to a portrait
report page. This produces a tall, large-element version (screenshots/flow-diagram.png)
that fills the report's text column and stays legible.

Run:    python scripts/build_flow_figure.py   (then office_export.py exports the PNG)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "HP-Bot-flow-figure.pptx"

W, H = Inches(7.5), Inches(8.0)
PAPER = RGBColor(0xFA, 0xF9, 0xF7)
INK = RGBColor(0x14, 0x14, 0x14)
MUTE = RGBColor(0x76, 0x72, 0x6B)
RULE = RGBColor(0xDD, 0xD8, 0xCF)
PANEL = RGBColor(0xF1, 0xEE, 0xE8)
GOLD = RGBColor(0xA0, 0x7D, 0x3D)
GOOD = RGBColor(0x44, 0x70, 0x4E)
BAD = RGBColor(0x9A, 0x3B, 0x3B)
FONT = "Segoe UI"
FONT_SEMI = "Segoe UI Semibold"
FONT_LIGHT = "Segoe UI Light"


def _run(p, text, *, font=FONT, size=14, color=INK, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def _box(s, x, y, w, h, title, sub, *, fill=RGBColor(0xFF, 0xFF, 0xFF), border=INK, tcol=INK, scol=MUTE, weight=1.4):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border; sp.line.width = Pt(weight)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
    _run(p, title, font=FONT_SEMI, size=15, color=tcol, bold=True)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = 1.0
        _run(p2, sub, font=FONT, size=10.5, color=scol)
    return sp


def _conn(s, x1, y1, x2, y2, color=INK, weight=1.5):
    ln = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False
    return ln


def _label(s, x, y, w, text, *, color=MUTE, size=10, align=PP_ALIGN.CENTER, italic=True):
    b = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.32))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    _run(p, text, size=size, color=color, italic=italic)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER; bg.line.fill.background(); bg.shadow.inherit = False

    sx, sw, h = 0.45, 3.75, 0.82      # spine x, width, box height
    cxr = sx + sw                      # right edge of spine
    ys = [0.95, 2.0, 3.05, 4.10, 5.15, 6.20]

    _box(s, sx, ys[0], sw, 0.66, "user message · Streamlit", "", fill=INK, border=INK, tcol=PAPER)
    _box(s, sx, ys[1], sw, h, "guard.py", "regex jailbreak prefilter")
    _box(s, sx, ys[2], sw, h, "embed query", "MiniLM-L6-v2")
    _box(s, sx, ys[3], sw, h, "Index A — questions only", "FAISS cosine · top-1")
    _box(s, sx, ys[4], sw, h, "Index B — all chunks", "hybrid 0.7 dense + 0.3 BM25 · top-5 → prompt")
    _box(s, sx, ys[5], sw, h, "OpenRouter chat completion", "temperature 0 · 7-model fallback")

    spine_x = sx + sw / 2
    seg = [(ys[0] + 0.66, ys[1]), (ys[1] + h, ys[2]), (ys[2] + h, ys[3]), (ys[3] + h, ys[4]), (ys[4] + h, ys[5])]
    for a, b in seg:
        _conn(s, spine_x, a, spine_x, b)

    # refuse branch (right of guard)
    bx = cxr + 0.45
    bw = 7.5 - bx - 0.25
    rf = _box(s, bx, ys[1], bw, h, "refuse", "no API call", fill=PANEL, border=BAD, tcol=BAD, scol=BAD)
    _conn(s, cxr, ys[1] + h / 2, bx, ys[1] + h / 2, color=BAD)
    _label(s, bx, ys[1] - 0.34, bw, "if jailbreak pattern", color=BAD)

    # cached branch (right of Index A)
    ca = _box(s, bx, ys[3], bw, h, "cached answer", "no LLM call", fill=PANEL, border=GOOD, tcol=GOOD, scol=GOOD)
    _conn(s, cxr, ys[3] + h / 2, bx, ys[3] + h / 2, color=GOOD)
    _label(s, bx, ys[3] - 0.34, bw, "if top-1 score ≥ 0.85", color=GOOD)

    # closing note
    b = s.shapes.add_textbox(Inches(0.45), Inches(ys[5] + h + 0.28), Inches(6.8), Inches(0.7))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT; p.line_spacing = 1.1
    _run(p, "The guard and the question cache resolve the easy cases instantly and for free; "
            "only genuinely new questions reach the model, and the answer is appended to memory.",
         size=11, color=MUTE, italic=True)

    prs.save(str(OUT))
    print(f"wrote {OUT.name}")


if __name__ == "__main__":
    build()
